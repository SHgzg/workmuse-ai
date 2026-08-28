from __future__ import annotations

import tempfile
import unittest
from pathlib import Path

from workmuse_worker.adapters import _transcript_blocks, _understand_with_python_libraries, frame_interval_seconds
from workmuse_worker.content import _is_relative_to, inspect_resource


class BuiltinAdapterTests(unittest.IsolatedAsyncioTestCase):
    async def asyncSetUp(self) -> None:
        self.events = []

    async def progress(self, event: str, data: dict) -> None:
        self.events.append((event, data))

    async def test_allowed_root_is_normalized_before_containment_check(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            nested = root / "nested"
            nested.mkdir()
            equivalent_root = nested / ".."
            self.assertTrue(_is_relative_to((root / "resource.txt").resolve(), equivalent_root))

    async def test_pdf_pages_are_extracted_with_page_evidence(self) -> None:
        from reportlab.pdfgen.canvas import Canvas

        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            path = root / "report.pdf"
            canvas = Canvas(str(path))
            canvas.drawString(72, 720, "Project roadmap review is due Friday")
            canvas.showPage()
            canvas.drawString(72, 720, "Testing time is the main delivery risk")
            canvas.save()
            result = await _understand_with_python_libraries(inspect_resource(path, (root,)), self.progress)
            self.assertEqual(result["metadata"]["pageCount"], 2)
            self.assertEqual(result["blocks"][0]["location"], {"kind": "page", "page": 1})

    async def test_docx_and_table_are_extracted(self) -> None:
        from docx import Document

        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            path = root / "meeting.docx"
            document = Document()
            document.add_heading("产品周会", level=1)
            document.add_paragraph("周五完成评审")
            table = document.add_table(rows=1, cols=2)
            table.cell(0, 0).text = "负责人"
            table.cell(0, 1).text = "陈一"
            document.save(path)
            result = await _understand_with_python_libraries(inspect_resource(path, (root,)), self.progress)
            self.assertIn("heading", [block["type"] for block in result["blocks"]])
            self.assertIn("table", [block["type"] for block in result["blocks"]])

    async def test_pptx_slides_keep_page_numbers(self) -> None:
        from pptx import Presentation

        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            path = root / "roadmap.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "路线图"
            slide.placeholders[1].text = "第一阶段完成资料理解"
            presentation.save(path)
            result = await _understand_with_python_libraries(inspect_resource(path, (root,)), self.progress)
            self.assertEqual(result["blocks"][0]["location"], {"kind": "page", "page": 1})

    async def test_image_metadata_is_available_without_ocr(self) -> None:
        from PIL import Image

        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary)
            path = root / "diagram.png"
            Image.new("RGB", (320, 180), "white").save(path)
            result = await _understand_with_python_libraries(inspect_resource(path, (root,)), self.progress)
            self.assertEqual(result["metadata"]["width"], 320)
            self.assertEqual(result["blocks"][0]["type"], "image")

    async def test_transcript_segments_keep_time_evidence(self) -> None:
        blocks = _transcript_blocks({
            "segments": [{"start": 1.25, "end": 3.5, "text": "确认产品路线图"}]
        })
        self.assertEqual(blocks[0]["location"], {"kind": "time", "startMs": 1250, "endMs": 3500})

    async def test_plain_transcript_remains_usable_without_segments(self) -> None:
        blocks = _transcript_blocks({"text": "会议决定周五完成评审"})
        self.assertEqual(blocks[0]["location"], {"kind": "resource"})

    async def test_video_frame_interval_is_bounded(self) -> None:
        self.assertEqual(frame_interval_seconds(30), 5)
        self.assertEqual(frame_interval_seconds(600), 60)
        self.assertEqual(frame_interval_seconds(10_000), 120)
