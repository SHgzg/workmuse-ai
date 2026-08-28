from __future__ import annotations

import asyncio
import json
import tempfile
from pathlib import Path
from typing import Any, Awaitable, Callable

from .content import Resource, normalized_content, provenance
from .tools import probe_tool, resolve_tool, sanitized_environment

Progress = Callable[[str, dict[str, Any]], Awaitable[None]]


async def understand_document(resource: Resource, progress: Progress, language: str | None = None) -> dict[str, Any]:
    errors: list[str] = []
    try:
        return await _understand_with_docling(resource, progress, language)
    except Exception as error:
        errors.append(f"Docling: {error}")

    fallback_id = "tesseract" if resource.kind == "image" else "mineru"
    try:
        return await (
            _understand_image_with_tesseract(resource, progress, language)
            if fallback_id == "tesseract"
            else _understand_document_with_mineru(resource, progress)
        )
    except Exception as error:
        errors.append(f"{fallback_id}: {error}")
    try:
        return await _understand_with_python_libraries(resource, progress)
    except Exception as error:
        errors.append(f"python-library: {error}")
    raise RuntimeError("No document adapter succeeded. " + " | ".join(errors))


async def _understand_with_docling(resource: Resource, progress: Progress, language: str | None) -> dict[str, Any]:
    probe = await probe_tool("docling")
    if not probe["available"]:
        raise FileNotFoundError("Docling is unavailable")

    _definition, executable = resolve_tool("docling")
    with tempfile.TemporaryDirectory(prefix="workmuse-docling-") as temporary:
        output = Path(temporary)
        args = ["convert", str(resource.path), "--to", "json", "--output", str(output)]
        if language:
            args.extend(["--ocr-lang", language])
        await progress("stage", {"name": "document.extract", "adapter": "docling"})
        await run_checked(executable, args, timeout=900)
        candidates = list(output.rglob("*.json"))
        if not candidates:
            raise ValueError("Docling completed without producing JSON")
        raw = json.loads(candidates[0].read_text(encoding="utf-8"))

    blocks = extract_docling_blocks(raw, resource.id)
    if not blocks:
        raise ValueError("Docling output did not contain recognizable content blocks")
    return normalized_content(
        resource,
        blocks,
        metadata={"sourceFormat": "docling-json"},
        provenance=[provenance("document-extract", "docling", str(probe.get("version", "unknown")))],
    )


async def _understand_document_with_mineru(resource: Resource, progress: Progress) -> dict[str, Any]:
    probe = await probe_tool("mineru")
    if not probe["available"]:
        raise FileNotFoundError("MinerU is unavailable")
    _definition, executable = resolve_tool("mineru")
    with tempfile.TemporaryDirectory(prefix="workmuse-mineru-") as temporary:
        output = Path(temporary)
        await progress("stage", {"name": "document.extract", "adapter": "mineru"})
        await run_checked(executable, ["-p", str(resource.path), "-o", str(output)], timeout=1800)
        markdown_files = list(output.rglob("*.md"))
        if not markdown_files:
            raise ValueError("MinerU completed without producing Markdown")
        text = markdown_files[0].read_text(encoding="utf-8", errors="replace")
    blocks = _plain_text_blocks(text)
    if not blocks:
        raise ValueError("MinerU output did not contain text")
    return normalized_content(
        resource,
        blocks,
        metadata={"sourceFormat": "mineru-markdown"},
        provenance=[provenance("document-extract", "mineru", str(probe.get("version", "unknown")))],
        warnings=[{"code": "document_fallback", "message": "Docling failed; MinerU fallback was used."}],
    )


async def _understand_image_with_tesseract(
    resource: Resource, progress: Progress, language: str | None
) -> dict[str, Any]:
    probe = await probe_tool("tesseract")
    if not probe["available"]:
        raise FileNotFoundError("Tesseract is unavailable")
    _definition, executable = resolve_tool("tesseract")
    args = [str(resource.path), "stdout"]
    if language:
        args.extend(["-l", language])
    await progress("stage", {"name": "image.ocr", "adapter": "tesseract"})
    text = await run_checked(executable, args, timeout=300)
    blocks = _plain_text_blocks(text)
    if not blocks:
        raise ValueError("Tesseract did not recognize text")
    return normalized_content(
        resource,
        blocks,
        metadata={"sourceFormat": "tesseract-text"},
        provenance=[provenance("image-ocr", "tesseract", str(probe.get("version", "unknown")))],
        warnings=[{"code": "image_fallback", "message": "Docling failed; Tesseract fallback was used."}],
    )


async def _understand_with_python_libraries(resource: Resource, progress: Progress) -> dict[str, Any]:
    await progress("stage", {"name": "document.extract", "adapter": "python-library"})
    extension = resource.path.suffix.lower()
    blocks: list[dict[str, Any]] = []
    metadata: dict[str, Any] = {"sourceFormat": "python-library"}

    if extension == ".pdf":
        from pypdf import PdfReader  # type: ignore[import-not-found]

        reader = PdfReader(str(resource.path))
        metadata["pageCount"] = len(reader.pages)
        for page_number, page in enumerate(reader.pages, start=1):
            text = (page.extract_text() or "").strip()
            if text:
                blocks.append(_block(blocks, text, {"kind": "page", "page": page_number}))
    elif extension in {".docx", ".doc"}:
        if extension == ".doc":
            raise ValueError("Legacy .doc requires Docling or LibreOffice conversion")
        from docx import Document  # type: ignore[import-not-found]

        document = Document(str(resource.path))
        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if text:
                block_type = "heading" if paragraph.style and paragraph.style.name.startswith("Heading") else "paragraph"
                blocks.append(_block(blocks, text, {"kind": "resource"}, block_type))
        for table in document.tables:
            rows = [" | ".join(cell.text.strip() for cell in row.cells) for row in table.rows]
            text = "\n".join(row for row in rows if row.strip(" |"))
            if text:
                blocks.append(_block(blocks, text, {"kind": "resource"}, "table"))
    elif extension in {".pptx", ".ppt"}:
        if extension == ".ppt":
            raise ValueError("Legacy .ppt requires Docling or LibreOffice conversion")
        from pptx import Presentation  # type: ignore[import-not-found]

        presentation = Presentation(str(resource.path))
        metadata["pageCount"] = len(presentation.slides)
        for slide_number, slide in enumerate(presentation.slides, start=1):
            texts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            if texts:
                blocks.append(_block(blocks, "\n".join(texts), {"kind": "page", "page": slide_number}))
    elif extension in {".xlsx", ".xls"}:
        if extension == ".xls":
            raise ValueError("Legacy .xls requires Docling or LibreOffice conversion")
        from openpyxl import load_workbook  # type: ignore[import-not-found]

        workbook = load_workbook(resource.path, read_only=True, data_only=False)
        metadata["sheets"] = workbook.sheetnames
        for sheet in workbook.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                values = ["" if value is None else str(value) for value in row]
                if any(values):
                    rows.append(" | ".join(values))
            if rows:
                blocks.append(_block(
                    blocks,
                    f"工作表：{sheet.title}\n" + "\n".join(rows),
                    {"kind": "resource"},
                    "table",
                ))
        workbook.close()
    elif resource.kind == "image":
        from PIL import Image  # type: ignore[import-not-found]

        with Image.open(resource.path) as image:
            metadata.update({"width": image.width, "height": image.height, "format": image.format, "mode": image.mode})
        blocks.append(_block(
            blocks,
            f"图片 {resource.file_name}，尺寸 {metadata['width']}×{metadata['height']}，格式 {metadata['format']}",
            {"kind": "resource"},
            "image",
        ))
    else:
        raise ValueError(f"No built-in library adapter for {extension}")

    if not blocks:
        raise ValueError("The built-in parser found no extractable content")
    return normalized_content(
        resource,
        blocks,
        metadata=metadata,
        provenance=[provenance("document-extract", "python-library")],
        warnings=[{
            "code": "library_fallback",
            "message": "External document tools were unavailable; a lightweight local parser was used.",
        }],
    )


async def understand_media(
    resource: Resource,
    progress: Progress,
    language: str | None = None,
    allow_cloud: bool = False,
) -> dict[str, Any]:
    try:
        return await _understand_media_local(resource, progress, language, allow_cloud)
    except Exception as local_error:
        from .ai import transcribe_file

        await progress("stage", {"name": "audio.transcribe", "adapter": "configured-provider"})
        transcript = await transcribe_file(str(resource.path), allow_cloud=allow_cloud, language=language)
        if transcript is None:
            raise RuntimeError(f"Local transcription failed and no permitted provider is configured: {local_error}")
        blocks = _transcript_blocks(transcript)
        if not blocks:
            raise ValueError("Transcription provider returned no text")
        return normalized_content(
            resource,
            blocks,
            metadata={"sourceFormat": "transcription-provider"},
            provenance=[provenance("audio-transcribe", "configured-provider")],
            warnings=[{"code": "transcription_fallback", "message": "A configured transcription provider was used."}],
        )


async def _understand_media_local(
    resource: Resource, progress: Progress, language: str | None, allow_cloud: bool
) -> dict[str, Any]:
    ffprobe = await probe_tool("ffprobe")
    whisper = await probe_tool("whisper")
    if not ffprobe["available"]:
        raise FileNotFoundError("Media understanding requires ffprobe")
    if not whisper["available"]:
        raise FileNotFoundError("Audio transcription requires Whisper")

    _probe_definition, probe_executable = resolve_tool("ffprobe")
    await progress("stage", {"name": "media.inspect", "adapter": "ffprobe"})
    metadata_output = await run_checked(
        probe_executable,
        ["-v", "error", "-show_format", "-show_streams", "-of", "json", str(resource.path)],
        timeout=60,
    )
    media_metadata = json.loads(metadata_output)

    _whisper_definition, whisper_executable = resolve_tool("whisper")
    with tempfile.TemporaryDirectory(prefix="workmuse-whisper-") as temporary:
        output = Path(temporary)
        args = [str(resource.path), "--output_format", "json", "--output_dir", str(output)]
        if language:
            args.extend(["--language", language])
        await progress("stage", {"name": "audio.transcribe", "adapter": "whisper"})
        await run_checked(whisper_executable, args, timeout=3600)
        candidates = list(output.glob("*.json"))
        if not candidates:
            raise ValueError("Whisper completed without producing JSON")
        transcript = json.loads(candidates[0].read_text(encoding="utf-8"))

    blocks = _transcript_blocks(transcript)

    content = normalized_content(
        resource,
        blocks,
        metadata={"media": media_metadata.get("format", {}), "streams": media_metadata.get("streams", [])},
        provenance=[
            provenance("media-inspect", "ffprobe", str(ffprobe.get("version", "unknown"))),
            provenance("audio-transcribe", "whisper", str(whisper.get("version", "unknown"))),
        ],
    )
    if resource.kind == "video":
        ffmpeg = await probe_tool("ffmpeg")
        if ffmpeg["available"]:
            from .ai import enrich_video_frames

            duration = float(media_metadata.get("format", {}).get("duration", 0) or 0)
            with tempfile.TemporaryDirectory(prefix="workmuse-frames-") as frame_directory:
                frames = await extract_video_frames(
                    resource.path,
                    Path(frame_directory),
                    frame_interval_seconds(duration),
                    progress,
                )
                content = await enrich_video_frames(content, frames, allow_cloud=allow_cloud)
        else:
            content["warnings"].append({
                "code": "video_frames_skipped",
                "message": "FFmpeg is unavailable; only the audio track was understood.",
            })
    return content


def _transcript_blocks(transcript: dict[str, Any]) -> list[dict[str, Any]]:
    blocks = []
    for index, segment in enumerate(transcript.get("segments", [])):
        text = str(segment.get("text", "")).strip()
        if text:
            blocks.append({
                "id": f"segment-{index + 1}",
                "type": "transcript",
                "text": text,
                "location": {
                    "kind": "time",
                    "startMs": round(float(segment.get("start", 0)) * 1000),
                    "endMs": round(float(segment.get("end", 0)) * 1000),
                },
            })
    if not blocks:
        text = str(transcript.get("text", "")).strip()
        if text:
            blocks.append({
                "id": "segment-1",
                "type": "transcript",
                "text": text,
                "location": {"kind": "resource"},
            })
    return blocks


def frame_interval_seconds(duration_seconds: float) -> int:
    if duration_seconds <= 0:
        return 30
    return max(5, min(120, round(duration_seconds / 10)))


async def extract_video_frames(
    path: Path, output: Path, interval_seconds: int, progress: Progress
) -> list[tuple[str, int]]:
    _definition, executable = resolve_tool("ffmpeg")
    await progress("stage", {"name": "video.keyframes", "adapter": "ffmpeg"})
    pattern = output / "frame-%03d.jpg"
    await run_checked(
        executable,
        ["-v", "error", "-i", str(path), "-vf", f"fps=1/{interval_seconds}", "-frames:v", "12", str(pattern)],
        timeout=600,
    )
    return [
        (str(frame), index * interval_seconds * 1000)
        for index, frame in enumerate(sorted(output.glob("frame-*.jpg")))
    ]


def extract_docling_blocks(raw: Any, resource_id: str) -> list[dict[str, Any]]:
    blocks: list[dict[str, Any]] = []
    seen: set[str] = set()

    def visit(value: Any) -> None:
        if isinstance(value, dict):
            text = value.get("text")
            if isinstance(text, str) and text.strip() and text not in seen:
                seen.add(text)
                block_type = _docling_type(value)
                block: dict[str, Any] = {
                    "id": f"block-{len(blocks) + 1}",
                    "type": block_type,
                    "text": text.strip(),
                    "location": _docling_location(value),
                }
                blocks.append(block)
            for child in value.values():
                visit(child)
        elif isinstance(value, list):
            for child in value:
                visit(child)

    visit(raw)
    return blocks


def _docling_type(value: dict[str, Any]) -> str:
    label = str(value.get("label", value.get("type", ""))).lower()
    if "title" in label or "heading" in label or "section" in label:
        return "heading"
    if "table" in label:
        return "table"
    if "list" in label:
        return "list"
    if "picture" in label or "image" in label:
        return "image"
    return "paragraph"


def _docling_location(value: dict[str, Any]) -> dict[str, Any]:
    provenance_items = value.get("prov")
    if isinstance(provenance_items, list) and provenance_items:
        item = provenance_items[0]
        if isinstance(item, dict):
            page = item.get("page_no", item.get("page", 1))
            location: dict[str, Any] = {"kind": "page", "page": int(page)}
            bbox = item.get("bbox")
            if isinstance(bbox, dict):
                coordinates = [bbox.get(key) for key in ("l", "t", "r", "b")]
                if all(isinstance(number, (int, float)) for number in coordinates):
                    location["bbox"] = coordinates
            return location
    return {"kind": "resource"}


def _plain_text_blocks(text: str) -> list[dict[str, Any]]:
    blocks = []
    for paragraph in (item.strip() for item in text.split("\n\n")):
        if paragraph:
            blocks.append({
                "id": f"block-{len(blocks) + 1}",
                "type": "paragraph",
                "text": paragraph,
                "location": {"kind": "resource"},
            })
    return blocks


def _block(
    blocks: list[dict[str, Any]], text: str, location: dict[str, Any], block_type: str = "paragraph"
) -> dict[str, Any]:
    return {
        "id": f"block-{len(blocks) + 1}",
        "type": block_type,
        "text": text,
        "location": location,
    }


async def run_checked(executable: str, args: list[str], timeout: float) -> str:
    process = await asyncio.create_subprocess_exec(
        executable,
        *args,
        stdout=asyncio.subprocess.PIPE,
        stderr=asyncio.subprocess.PIPE,
        env=sanitized_environment(),
    )
    try:
        stdout, stderr = await asyncio.wait_for(process.communicate(), timeout=timeout)
    except (asyncio.TimeoutError, asyncio.CancelledError):
        process.kill()
        await process.wait()
        raise
    if process.returncode != 0:
        detail = stderr.decode("utf-8", errors="replace")[-4000:]
        raise RuntimeError(f"Tool failed with exit code {process.returncode}: {detail}")
    return stdout.decode("utf-8", errors="replace")
